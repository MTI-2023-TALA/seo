from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE
from bs4 import BeautifulSoup as bs
from flask import Flask, send_file, request, jsonify
import requests
import re
import validators

app = Flask(__name__)

class Heading:
  def __init__(self, name, text):
    self.name = name
    self.text = text
  
  def get_name(self):
    return self.name

  def get_text(self):
    return self.text
  
  def nb_spaces(self):
    if self.name == "h2":
      return "\t"
    if self.name == "h3":
      return "\t\t"
    if self.name == "h4":
      return "\t\t\t"
    if self.name == "h5":
      return "\t\t\t\t"
    if self.name == "h6":
      return "\t\t\t\t\t"
    else:
      return ""

  def __str__(self):
    return self.name + ": " + self.text

@app.route('/', methods=['GET'])
def index():
  return send_file('static/index.html')


@app.route('/', methods=['POST'])
def hello_world():
  url = request.json['url']
  valid = validators.url(url)
  if not valid:
    print('toto?')
    raise Exception("Please provide a valid URL to analyze!")

  response = requests.get(url)
  html = response.content
  soup = bs(html, "lxml")

  prs=Presentation()
  lyt=prs.slide_layouts[6]
  left = 0
  top = 1
  width = Inches(10)
  height = Inches(7)


  def add_begin_page_ppt():
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.add_paragraph()
    p.text = "SEO Report" 
    p.font.size = Pt(30)
    p.font.bold = True

    p = tf.add_paragraph() 
    p.text = "From: " + url
    p.font.bold = True


  def get_title(soup):
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    balise_title = soup.find('title').string
    title_length = len(balise_title)
    title_words = len(balise_title.split(' '))

    p = tf.add_paragraph()
    p.text = "Contenu de la balise titre : \"" + balise_title + "\"\n"
    p.font.bold = True

    p = tf.add_paragraph() 
    p.text = "Le titre de la page contient : " + str(title_length) + " caractères et " + str(title_words) + " mots."

    p = tf.add_paragraph()
    p.font.italic = True
    if title_words < 7:
      p.text = "Le titre est trop court, pensez à le rallonger pour atteindre 10 à 12 mots."
    elif title_words < 10:
      p.text = "Le titre est légèrement trop court, pensez à le rallonger un petit peu pour atteindre 10 à 12 mots."
    else:
      p.text = "Votre titre a une bon nombre de mots, c'est bien !"
    
    p = tf.add_paragraph() 
    p.text = "CONSEIL: Un bon titre contient entre 10 et 12 mots et ne dépasse pas 70 caractères !\n" + \
            "CONSEIL: Proposez un titre différent et cohérent pour chacunes de vos pages.\n" + \
            "CONSEIL: Evitez de répeter deux fois un même mot dans votre titre."
    p.font.italic = True


  def get_meta_description(soup):
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    description = soup.find("meta", property="og:description")
    if description:
      description_text = description.get("content")
      description_length = len(description_text)
      description_words = len(description_text.split(' '))
      
      p = tf.add_paragraph()
      p.text = "Contenu de la balise méta description : \"" + description_text + "\"\n"
      p.font.bold = True

      p = tf.add_paragraph()
      p.text = "La balise méta description contient : " + str(description_length) + " caractères et " + str(description_words) + " mots.\n"

      p = tf.add_paragraph()
      p.font.italic = True
      if description_length < 200:
        p.text = "Pensez à ajouter un peu de texte dans la balise méta description pour atteindre 200 à 300 caractères (espaces compris)."
      elif description_length > 300:
        p.text = "Pensez à retirer un peu de texte dans la balise méta description pour être entre 200 à 300 caractères (espaces compris). Soyez plus conçis et accrocheur !"
      else:
        p.text = "Le nombre de caractères de votre méta description correspond bien dans les limites"

    else:
      p = tf.add_paragraph()
      p.text = "La balise méta description n'est pas présente.\n"
      p.font.bold = True

      p = tf.add_paragraph()
      p.text = "Il est important d'en renseigner une afin que votre site soit bien référencé!\nPensez à y mettre une phrase accrocheuse et ne séparez pas les mots par des virgules."
    
    p = tf.add_paragraph() 
    p.text = "CONSEIL: Pensez à placer les mots clés importants dans les 150 premiers caractères de votre description afin d'être mieux référencé.\n" + \
            "CONSEIL: Pensez à ajouter une méta-description différente sur chacune de vos pages et qu'elle soit assez longue."
    p.font.italic = True


  def get_meta_keywords(soup):
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    keywords = soup.find("meta", property="og:keywords")
    if keywords:
      keywords_text = keywords.get("content")

      p = tf.add_paragraph()
      p.text = "Contenu de la balise méta Keywords : \"" + keywords_text + "\"\n"
      p.font.bold = True

      p = tf.add_paragraph()
      p.text = "Ne tenez pas compte de cette balise! Google ne la lit plus et les concurrents de Google ne la lisent pas ou alors lui donnent une importance minime !"
    else:
      p = tf.add_paragraph()
      p.text = "La balise méta Keywords n'est pas présente.\n"
      p.font.bold = True

      p = tf.add_paragraph()
      p.font.italic = True
      p.text = "Ne vous en faites pas, ne tenez pas compte de cette balise! Google ne la lit plus et les concurrents de Google ne la lisent pas ou alors lui donnent une importance minime !"


  def get_url(soup):

    def verif_url(s):
      if '_' in s:
        return False
      return all(ord(c) < 128 for c in s)
    
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    url = soup.find("meta", property="og:url")
    if url:
      url_text = url.get("content")
      url_length = len(url_text)
      
      p = tf.add_paragraph()
      p.text = "URL donnée : \"" + url_text + "\"\n"
      p.font.bold = True

      p = tf.add_paragraph()
      p.text = "L'URL fait " + str(url_length) + " caractères.\n"

      if not verif_url(url_text):
        p = tf.add_paragraph()
        p.text = "Attention votre URL contient des underscores ou des caractères accentués ! Ce n'est pas une bonne chose, pensez à corriger ceci."
        p.font.bold = True
      else:
        p = tf.add_paragraph()
        p.font.bold = True
        p.text = "Votre URL ne contient ni underscore ni caractères accentués ! C'est très bien.\n"

    else:
      p = tf.add_paragraph()
      p.text = "Aucune URL n'a été trouvée.\n"
      p.font.bold = True
    
    p = tf.add_paragraph() 
    p.text = "CONSEIL: Il faut qu'à la lecture de l'URL nous puissions comprendre ce que le contenu de la page va être!\n" + \
              "CONSEIL: Evitez les \"?\" et \"&\" dans l'URL.\n" + \
              "CONSEIL: Si vous devez séparer des mots, utilisez le tiret du haut \"-\".\n" + \
              "CONSEIL: N'utilisez pas de mots accentués dans l'URL."
    p.font.italic = True

  def get_structure(soup):
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    headings = soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6"])

    formatted_headings = []
    for heading in headings:
      formatted_headings.append(Heading(heading.name, heading.text.strip()))
    
    p = tf.add_paragraph()
    p.text = "Structure Hn"
    p.font.bold = True

    count = 1
    for h in formatted_headings:
      if count % 21 == 0:
        slide=prs.slides.add_slide(lyt)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True

        p = tf.add_paragraph()
        p.text = "Suite Structure Hn"
        p.font.bold = True
      p = tf.add_paragraph()
      p.text = h.nb_spaces() + "[" + h.name + "] " + h.text
      count += 1
    
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "Conseils sur la structure Hn"
    p.font.bold = True
    p = tf.add_paragraph() 
    p.text = "La structure affichée représente l'arborescence de votre page, il est nécessaire qu'à la lecture de celle-ci " + \
              "nous puissons remarquer que la page a du sens et est bien architecturée.\n" + \
              "Il n'est pas grave de sauter des niveaux de Hn : exemple passer de H3 à H6.\n" + \
              "Il est possible de mettre plusieurs balises H1 dans une page mais en général nous préférons n'en garder qu'une."
    p.font.italic = True

  def get_textual_content(soup):
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    d = {}
    decoded_html = html.decode("utf-8")
    decoded_html_split = re.split('\\s+', decoded_html)

    for word in decoded_html_split:
      if word in d:
        d[word] += 1
      else:
        d[word] = 1

    sorted_most_used_words = []
    for k, v in sorted(d.items(), key=lambda kv: kv[1], reverse=True):
      sorted_most_used_words.append("%s ➡ %s occurences" % (k,v))
    most_used_words = sorted_most_used_words[:10]

    p = tf.add_paragraph()
    p.text = "Mots clés les plus utilisés :\n"
    p.font.bold = True

    for word in most_used_words:
      p = tf.add_paragraph()
      p.text = word


  def get_other_content(soup):
    slide=prs.slides.add_slide(lyt)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "Autres contenus\n"
    p.font.bold = True

    robots = soup.find("meta", property="og:robots")
    p = tf.add_paragraph()
    if robots:
      robots_text = robots.get("content")
      p.text = "La contenu de la balise Robots est le suivant : \"" + robots_text + "\""
    else:
      p.text = "La balise méta Robots n'est pas présente."


    canonicals = soup.select('link[rel*=canonical]')
    p = tf.add_paragraph()
    if canonicals:
      canonicals_text = '\n'.join(canonical['href'] for canonical in canonicals)
      p.text = "La balise canonical est la suivante : \"" + canonicals_text + "\""
    else:
      p.text = "La balise canonical n'est pas présente."
    
    list_hreflangs = [[a['href'], a["hreflang"]] for a in soup.find_all('link', href=True, hreflang=True)]
    p = tf.add_paragraph()
    if list_hreflangs:
      p.text = "Les balises hreflang sont les suivantes : \"" + list_hreflangs + "\""
    else:
      p.text = "Les balises hreflang ne sont pas présentes."

    nb_images = soup.find_all('img')
    nb_images_text = len(nb_images)
    p = tf.add_paragraph()
    p.text = "Il y a " + str(nb_images_text) + " images dans la page."

    nb_images_alt = soup.find_all('img', alt=True)
    nb_images_alt_text = len(nb_images_alt)
    p = tf.add_paragraph()
    p.text = "Il y a " + str(nb_images_alt_text) + " images avec un attribut ALT rempli."

    nb_images_no_alt_text = nb_images_text - nb_images_alt_text
    p = tf.add_paragraph()
    p.text = "Il y a " + str(nb_images_no_alt_text) + " images avec un attribut ALT vide ou non présent."

    outbound_links = [[a.get_text(), a["href"]] for a in soup.find_all('a', href=True)]
    outbound_links_text = len(outbound_links)
    p = tf.add_paragraph()
    p.text = "Il y a " + str(outbound_links_text) + " liens sortants dans la page."

    p = tf.add_paragraph() 
    p.text = "CONSEIL: La balise méta Robots indique aux moteurs de recherche ce qu'ils doivent faire dans la page\n" + \
              "Détails ici -> http://robots-txt.com/meta-robots/" + \
              "CONSEIL: La balise Canonical permet d'éviter que si plusieurs pages ont le même contenu sur votre site web que cette page soit"  + \
              "explorée en priorité et que les pages \"doubles\" soient explorées moins souvent."
    p.font.italic = True

  add_begin_page_ppt()
  get_title(soup)
  get_meta_description(soup)
  get_meta_keywords(soup)
  get_url(soup)
  get_structure(soup)
  get_textual_content(soup)
  get_other_content(soup)
  prs.save("seo_result.pptx")
  return send_file("seo_result.pptx")
